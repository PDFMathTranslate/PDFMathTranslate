import os
import shutil
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Pt, RGBColor
from pptx.enum.text import PP_ALIGN
from pdf2zh.translator import BaseTranslator
from pdf2zh.config import ConfigManager
import logging

logger = logging.getLogger(__name__)

def translate_ppt(
    file_path: str,
    lang_from: str,
    lang_to: str,
    translator: BaseTranslator,
    output_dir: str,
    callback: Any = None,
    cancellation_event: Any = None,
    **kwargs
) -> str:
    """
    Translate a PowerPoint file from one language to another, preserving formatting.

    Inputs:
        - file_path: Path to the input PowerPoint file
        - lang_from: Source language code
        - lang_to: Target language code
        - translator: Translator instance to use
        - output_dir: Directory to save the translated PowerPoint file
        - callback: Callback function for progress updates
        - cancellation_event: Event to signal cancellation
        - kwargs: Additional parameters

    Returns:
        - Path to the translated PowerPoint file
    """
    # Load the PowerPoint presentation
    prs = Presentation(file_path)
    
    # Get all text from the presentation
    text_items = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                if not paragraph.text.strip():
                    continue
                text_items.append({
                    'slide_idx': slide_idx,
                    'shape': shape,
                    'para_idx': para_idx,
                    'text': paragraph.text
                })
    
    # Translate all text items
    if text_items:
        texts = [item['text'] for item in text_items]
        translated_texts = translator.translate_batch(texts, lang_from, lang_to)
        
        # Update the presentation with translated text
        for item, translated_text in zip(text_items, translated_texts):
            slide_idx = item['slide_idx']
            shape = item['shape']
            para_idx = item['para_idx']
            
            # Preserve the original formatting
            original_paragraph = shape.text_frame.paragraphs[para_idx]
            original_font = original_paragraph.runs[0].font if original_paragraph.runs else None
            
            # Clear the existing text
            original_paragraph.clear()
            
            # Add the translated text with preserved formatting
            run = original_paragraph.add_run()
            run.text = translated_text
            
            if original_font:
                run.font.name = original_font.name
                run.font.size = original_font.size
                run.font.bold = original_font.bold
                run.font.italic = original_font.italic
                run.font.underline = original_font.underline
                run.font.color.rgb = original_font.color.rgb
                run.font.superscript = original_font.superscript
                run.font.subscript = original_font.subscript
    
    # Save the translated presentation
    filename = os.path.splitext(os.path.basename(file_path))[0]
    output_file = os.path.join(output_dir, f"{filename}-translated.pptx")
    prs.save(output_file)
    
    return output_file

def process_ppt_file(
    file_path: str,
    lang_from: str,
    lang_to: str,
    service: str,
    output_dir: str,
    callback: Any = None,
    cancellation_event: Any = None,
    envs: Dict = None,
    prompt: Any = None,
    ignore_cache: bool = False,
    **kwargs
) -> str:
    """
    Process a PowerPoint file for translation.

    Inputs:
        - file_path: Path to the input PowerPoint file
        - lang_from: Source language code
        - lang_to: Target language code
        - service: Translation service to use
        - output_dir: Directory to save the translated file
        - callback: Callback function for progress updates
        - cancellation_event: Event to signal cancellation
        - envs: Environment variables for the translator
        - prompt: Custom prompt for the translator
        - ignore_cache: Whether to ignore the translation cache
        - kwargs: Additional parameters

    Returns:
        - Path to the translated PowerPoint file
    """
    from pdf2zh.translator import service_map
    
    # Create translator instance
    translator_class = service_map[service]
    translator = translator_class(
        lang_from,
        lang_to,
        "",
        envs=envs,
        prompt=prompt,
        ignore_cache=ignore_cache,
    )
    
    # Translate the PowerPoint file
    translated_file = translate_ppt(
        file_path,
        lang_from,
        lang_to,
        translator,
        output_dir,
        callback,
        cancellation_event,
        **kwargs
    )
    
    return translated_file